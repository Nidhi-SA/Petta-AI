import os
import pdfplumber
import pytesseract
from PIL import Image
import docx
from pptx import Presentation
import xlrd
import openpyxl

from sentence_transformers import SentenceTransformer
import faiss

# Load embedding model
model = SentenceTransformer("all-MiniLM-L6-v2")

DATA_DIR = [
    r"C:\Users\nidhi\Petta-AI\data",
]

# ---------------------------------------------------------
# FILE READERS
# ---------------------------------------------------------

def read_pdf(path):
    text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
    return text

def read_image(path):
    img = Image.open(path)
    return pytesseract.image_to_string(img)

def read_docx(path):
    try:
        doc = docx.Document(path)
        return [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    except:
        return []

def read_txt(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def read_pptx(path):
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        if texts:
            slides.append("\n".join(texts))
    return slides

def read_xlsx(path):
    ext = os.path.splitext(path)[1].lower()
    rows = []

    if ext == ".xls":
        wb = xlrd.open_workbook(path)
        for sheet in wb.sheets():
            for r in range(sheet.nrows):
                row = [str(sheet.cell(r, c).value) for c in range(sheet.ncols)]
                rows.append(" | ".join(row))
    else:
        wb = openpyxl.load_workbook(path)
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_vals = [str(cell) for cell in row]
                rows.append(" | ".join(row_vals))

    return rows

# ---------------------------------------------------------
# LOAD ALL DOCUMENTS
# ---------------------------------------------------------

def load_all_documents():
    texts = []
    filenames = []

    for folder in DATA_DIR:
        if not os.path.exists(folder):
            continue

        for root, _, files in os.walk(folder):
            for file in files:
                path = os.path.join(root, file)
                ext = file.lower()

                if ext.endswith(".pdf"):
                    texts.append(read_pdf(path))
                    filenames.append(path)

                elif ext.endswith((".png", ".jpg", ".jpeg")):
                    texts.append(read_image(path))
                    filenames.append(path)

                elif ext.endswith(".docx"):
                    for para in read_docx(path):
                        texts.append(para)
                        filenames.append(f"{path} (paragraph)")

                elif ext.endswith(".txt"):
                    texts.append(read_txt(path))
                    filenames.append(path)

                elif ext.endswith((".ppt", ".pptx")):
                    for slide in read_pptx(path):
                        texts.append(slide)
                        filenames.append(f"{path} (slide)")

                elif ext.endswith((".xls", ".xlsx")):
                    for row in read_xlsx(path):
                        texts.append(row)
                        filenames.append(f"{path} (row)")

    return filenames, texts

# ---------------------------------------------------------
# BUILD FAISS INDEX
# ---------------------------------------------------------

def build_faiss_index(texts):
    embeddings = model.encode(texts)
    dim = embeddings.shape[1]
    index = faiss.IndexFlatL2(dim)
    index.add(embeddings)
    return index

# ---------------------------------------------------------
# HYBRID SEARCH (Keyword + FAISS + Ranking)
# ---------------------------------------------------------

def hybrid_search(query, filenames, texts, index):
    keyword = query.lower()
    keyword_hits = []

    # Keyword search
    for filename, text in zip(filenames, texts):
        lines = text.split("\n")
        for line in lines:
            if keyword in line.lower():
                keyword_hits.append(f"[{filename}] {line.strip()}")

    # If keyword matches exist → return them ranked
    if keyword_hits:
        return "\n".join(keyword_hits)

    # Otherwise FAISS semantic search
    query_emb = model.encode([query])
    D, I = index.search(query_emb, k=5)

    results = []
    for idx in I[0]:
        results.append(f"[{filenames[idx]}]\n{texts[idx]}\n")

    return "\n".join(results)
